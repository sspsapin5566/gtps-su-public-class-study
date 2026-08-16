import pptx

prs = pptx.Presentation("output/slides.pptx")
print(f"Presentation has {len(prs.slides)} slides.")

for i, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            print(f"  Shape {j} (Text): {text[:100]}...")
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            print(f"  Shape {j} (Picture)")
