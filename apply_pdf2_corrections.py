import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def apply_pdf2_fixes():
    pptx_path = "output/slides.pptx"
    prs = pptx.Presentation(pptx_path)

    # -------------------------------------------------------------
    # 1. Slide 4 (Ppt 第四頁): 刪除圖片，使觀課時間序列紀錄全部排進 PPT
    # -------------------------------------------------------------
    slide4 = prs.slides[3] # 0-indexed slide 4
    # Delete picture shape on slide 4
    for shape in list(slide4.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            sp = shape._element
            sp.getparent().remove(sp)

    # Expand text box on slide 4 to full width
    for shape in slide4.shapes:
        if shape.has_text_frame:
            shape.left = Inches(0.8)
            shape.top = Inches(1.6)
            shape.width = Inches(11.7)
            shape.height = Inches(5.2)
            tf = shape.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.font.size = Pt(16)
                p.space_after = Pt(10)

    # -------------------------------------------------------------
    # 2. Slide 19, 20 (Ppt 第 19, 20 頁): 表格文字加大到 24 號字體
    # -------------------------------------------------------------
    for slide_idx in [18, 19]:
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Increase height and positioning to fit 24pt text
                shape.left = Inches(0.6)
                shape.top = Inches(1.6)
                shape.width = Inches(12.1)
                
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(24) # 24 號字體

    # -------------------------------------------------------------
    # 3. Slide 21 (Ppt 第 21 頁): 標題加大到 36 號字體，內文加大到 20 號字體
    # -------------------------------------------------------------
    slide21 = prs.slides[20] # Slide 21
    for shape in slide21.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            if "七、" in text or "結論" in text:
                for p in tf.paragraphs:
                    p.font.size = Pt(36) # 標題 36 號字
                    p.font.bold = True
            else:
                for p in tf.paragraphs:
                    if p.font.bold or "1." in p.text or "2." in p.text or "3." in p.text:
                        p.font.size = Pt(22) # 子標題 22 號字
                    else:
                        p.font.size = Pt(20) # 內文 20 號字

    # -------------------------------------------------------------
    # 4. Slide 2, 6, 10, 14, 18: 清除末尾空白無文字的黑點 (Bullet point)
    # -------------------------------------------------------------
    target_slides = [1, 5, 9, 13, 17] # 0-indexed for 2, 6, 10, 14, 18
    for s_idx in target_slides:
        slide = prs.slides[s_idx]
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                # Iterate through paragraphs and remove empty bullet paragraphs
                p_elements_to_remove = []
                for p in tf.paragraphs:
                    cleaned_txt = p.text.replace("•", "").replace("•", "").strip()
                    if not cleaned_txt:
                        p_elements_to_remove.append(p)
                
                for p in p_elements_to_remove:
                    # Remove paragraph element if empty
                    p_elem = p._p
                    p_elem.getparent().remove(p_elem)

    prs.save(pptx_path)
    print(f"Successfully applied PDF 2 fixes to {pptx_path}")

if __name__ == "__main__":
    apply_pdf2_fixes()
