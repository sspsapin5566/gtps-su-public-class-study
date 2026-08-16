import pptx

prs = pptx.Presentation("output/slides.pptx")

# 1. Slide 4 check
slide4 = prs.slides[3]
pics_s4 = [s for s in slide4.shapes if s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE]
print(f"Slide 4 pictures count: {len(pics_s4)} (Should be 0)")

# 2. Slide 19 & 20 table font size check
for idx in [18, 19]:
    slide = prs.slides[idx]
    for s in slide.shapes:
        if s.has_table:
            table = s.table
            p = table.cell(0, 0).text_frame.paragraphs[0]
            print(f"Slide {idx+1} table font size: {p.font.size.pt} pt (Should be 24 pt)")

# 3. Slide 21 font size check
slide21 = prs.slides[20]
for s in slide21.shapes:
    if s.has_text_frame:
        text = s.text_frame.text.strip()
        if "七、" in text:
            print(f"Slide 21 title font size: {s.text_frame.paragraphs[0].font.size.pt} pt (Should be 36 pt)")

# 4. Empty bullet check on Slides 2, 6, 10, 14, 18
for idx in [1, 5, 9, 13, 17]:
    slide = prs.slides[idx]
    empty_bullets = 0
    for s in slide.shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                cleaned = p.text.replace("•", "").strip()
                if not cleaned:
                    empty_bullets += 1
    print(f"Slide {idx+1} empty trailing bullets count: {empty_bullets} (Should be 0)")
