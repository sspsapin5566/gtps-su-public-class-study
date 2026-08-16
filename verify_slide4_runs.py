import pptx

prs = pptx.Presentation("output/slides.pptx")
slide4 = prs.slides[3]

print("=== SLIDE 4 NEW FORMATTING ===")
for shape in slide4.shapes:
    if shape.has_text_frame:
        print(f"Shape: {shape.name}")
        for i, p in enumerate(shape.text_frame.paragraphs):
            print(f"  P{i} [text={p.text[:70]}...]:")
            for r in p.runs:
                r_size = r.font.size.pt if r.font.size else "Inherit"
                print(f"    Run [size={r_size}, bold={r.font.bold}]: {r.text[:60]}")
