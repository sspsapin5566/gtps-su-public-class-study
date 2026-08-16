import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_perfect_presentation():
    old_pptx = "output/slides.pptx"
    out_pptx = "output/slides.pptx"
    
    # Images paths
    img_slide2 = "/Users/zhangwenbin/.gemini/antigravity/brain/73a5c8c4-3a1f-48b6-a788-2e684e569bc1/slide2_taiwanese_children_illustration_1786866270345.jpg"
    img_g1 = "/Users/zhangwenbin/.gemini/antigravity/brain/73a5c8c4-3a1f-48b6-a788-2e684e569bc1/group1_illustration_1786281760396.jpg"
    img_g2 = "/Users/zhangwenbin/.gemini/antigravity/brain/73a5c8c4-3a1f-48b6-a788-2e684e569bc1/group2_illustration_1786281775040.jpg"
    img_g3 = "/Users/zhangwenbin/.gemini/antigravity/brain/73a5c8c4-3a1f-48b6-a788-2e684e569bc1/group3_illustration_1786283728097.jpg"
    img_g4 = "/Users/zhangwenbin/.gemini/antigravity/brain/73a5c8c4-3a1f-48b6-a788-2e684e569bc1/group4_illustration_1786283742578.jpg"

    prs = pptx.Presentation(old_pptx)
    
    # 1. Global Replace "Study of Learning Community" -> "School as Learning Community, SLC"
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if "Study of Learning Community" in p.text:
                        p.text = p.text.replace("Study of Learning Community", "School as Learning Community, SLC")
                    if "Study of Learning Community, SLC" in p.text:
                        p.text = p.text.replace("Study of Learning Community, SLC", "School as Learning Community, SLC")

    # 2. Slide 2: Update picture with Taiwanese children skin tones illustration
    slide2 = prs.slides[1] # 0-indexed slide 2
    # Find picture shape on slide 2
    for shape in list(slide2.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            # Get position
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # Remove old shape
            sp = shape._element
            sp.getparent().remove(sp)
            # Add new picture
            slide2.shapes.add_picture(img_slide2, left, top, width, height)

    # 3. Slide 4 (Group 1 Description): Set picture match Slide 3 (Group 1 illustration)
    slide4 = prs.slides[3]
    for shape in list(slide4.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide4.shapes.add_picture(img_g1, left, top, width, height)

    # 4. Slide 7 (Group 2 Overview): Insert Group 2 Illustration Picture
    slide7 = prs.slides[6]
    has_pic7 = any(s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE for s in slide7.shapes)
    if not has_pic7:
        slide7.shapes.add_picture(img_g2, Inches(8.2), Inches(1.8), Inches(4.8), Inches(4.8))

    # 5. Slide 11 (Group 3 Overview): Insert Group 3 Illustration Picture
    slide11 = prs.slides[10]
    has_pic11 = any(s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE for s in slide11.shapes)
    if not has_pic11:
        slide11.shapes.add_picture(img_g3, Inches(8.2), Inches(1.8), Inches(4.8), Inches(4.8))

    # 6. Slide 15 (Group 4 Overview): Insert Group 4 Illustration Picture
    slide15 = prs.slides[14]
    has_pic15 = any(s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE for s in slide15.shapes)
    if not has_pic15:
        slide15.shapes.add_picture(img_g4, Inches(8.2), Inches(1.8), Inches(4.8), Inches(4.8))

    # 7. Slide 19 & Slide 20 (Matrix I & II): Replace text bullet markdown table with PPTX Table shapes
    # Slide 19 Data
    headers_19 = ["分析維度", "第一組 (Group 1)", "第二組 (Group 2)", "第三組 (Group 3)", "第四組 (Group 4)"]
    rows_19 = [
        ["成員組合與特性", "2男2女，有明顯沉思者與主筆者分工", "2男2女，成熟同儕傾聽與微鷹架輔導", "2男2女，對等雙點交替主導溝通", "1男3女，高密度多人同時主筆創作"],
        ["座位與空間動態", "圍坐討論，主筆居右前側", "前傾對焦，傾聽姿態顯著", "圍坐圓弧，流暢對等互動", "全員前傾，三人同時筆錄海報"],
        ["觀課互動特徵", "沉思者翻閱文本，主筆者分享記錄", "長髮女生引導，藍衣男生安心傾聽", "黃衣女與藍外套男生雙向交織", "低姿態陪伴，多點同步繪圖與紀錄"]
    ]

    # Slide 20 Data
    headers_20 = ["分析維度", "第一組 (Group 1)", "第二組 (Group 2)", "第三組 (Group 3)", "第四組 (Group 4)"]
    rows_20 = [
        ["傾聽品質", "沉思型傾聽（托腮關注）", "互動型傾聽（眼神對焦與前傾）", "平衡型傾聽（流暢輪流發言）", "全員前傾對焦（高度集中注意力）"],
        ["思考深度", "事實整理層（複製文本與畫圖）", "因果推論層（探究歷史發展背景）", "文本地圖整合層（雙向對照古今）", "個案推論與全班發表整合層"],
        ["跳躍伸展建議", "引導「為什麼」因果思考推論", "進行跨時空昔今系統化對照", "拋出衝突性歷史情境立論辯證", "海報加入心智圖網路與邏輯箭頭"]
    ]

    def create_table_on_slide(slide, headers, rows):
        # Remove text box containing old markdown table
        for shape in list(slide.shapes):
            if shape.has_text_frame and ("|" in shape.text_frame.text or "分析維度" in shape.text_frame.text):
                sp = shape._element
                sp.getparent().remove(sp)
                
        # Create PPTX Table
        x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8)
        shape_table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, cx, cy)
        table = shape_table.table

        # Set Column Widths
        table.columns[0].width = Inches(2.2) # Header col
        table.columns[1].width = Inches(2.37)
        table.columns[2].width = Inches(2.37)
        table.columns[3].width = Inches(2.37)
        table.columns[4].width = Inches(2.37)

        # Header formatting
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(29, 78, 216) # Rich Blue Header
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Rows formatting
        for row_idx, row_data in enumerate(rows):
            bg_color = RGBColor(248, 250, 252) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            for col_idx, text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(15, 23, 42)
                if col_idx == 0:
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                    cell.fill.fore_color.rgb = RGBColor(241, 245, 249)

    slide19 = prs.slides[18]
    create_table_on_slide(slide19, headers_19, rows_19)

    slide20 = prs.slides[19]
    create_table_on_slide(slide20, headers_20, rows_20)

    # 8. Delete Slide 21 ("視覺配色與簡報輸出說明")
    # In python-pptx, deleting slide requires removing rId element
    rId = prs.slides._sldIdLst[20].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[20]

    # 9. Add New Slide 21: Conclusion Slide (七、第 1 至第 4 小組 SLC 綜合分析結論與教學建議)
    blank_slide_layout = prs.slide_layouts[6]
    slide_conc = prs.slides.add_slide(blank_slide_layout)

    # Background color fill
    bg = slide_conc.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(253, 251, 247)

    # Title shape
    txBox = slide_conc.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 七、第 1 至第 4 小組 SLC 綜合分析結論與教學建議"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(74, 59, 50)

    # Content cards (3 Pillars)
    pillars = [
        {
            "title": "1. 傾聽與同儕鷹架文化 (Listening & Scaffolding)",
            "items": [
                "• 四個小組均展現高素質的互惠傾聽，第一組沉思型、第二組前傾對焦、第三組對等交織、第四組多點主筆。",
                "• 學生在同儕陪伴下表現安心發言，展現「School as Learning Community, SLC」互惠共學的核心精神。"
            ]
        },
        {
            "title": "2. 伸展跳躍課題與認知轉折 (Jump Learning & Cognitive Extension)",
            "items": [
                "• 各組從事實記錄成功邁向因果推論與古今對照（塭仔圳發展與新莊鐵路個案）。",
                "• 建議未來可拋出衝突性歷史情境題，引導小組進行高階觀點立論與邏輯辯證。"
            ]
        },
        {
            "title": "3. 教師角色與課堂陪伴策略 (Teacher's Role & Low-Profile Presence)",
            "items": [
                "• 蘇國瑞老師以彎腰低姿態傾聽、適切點撥問題鷹架，有效促進學生自主探究。",
                "• 建議巡視時可針對思考停滯小組拋出伸展跳躍問題，進一步推動深度跳躍學習。"
            ]
        }
    ]

    for idx, card in enumerate(pillars):
        top_y = Inches(1.8 + idx * 1.7)
        # Card Box
        shape_box = slide_conc.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y, Inches(11.7), Inches(1.5))
        shape_box.fill.solid()
        shape_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape_box.line.color.rgb = RGBColor(226, 232, 240)
        shape_box.line.width = Pt(1.5)

        tf_card = shape_box.text_frame
        tf_card.word_wrap = True
        
        # Header
        p0 = tf_card.paragraphs[0]
        p0.text = card["title"]
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(29, 78, 216)

        # Body items
        for it in card["items"]:
            p_item = tf_card.add_paragraph()
            p_item.text = it
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = RGBColor(51, 65, 85)

    prs.save(out_pptx)
    print(f"Saved PERFECT OVERHAULED presentation to {out_pptx}")

if __name__ == "__main__":
    build_perfect_presentation()
