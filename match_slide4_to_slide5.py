import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def update_slide4_formatting():
    pptx_path = "output/slides.pptx"
    prs = pptx.Presentation(pptx_path)

    slide4 = prs.slides[3] # 0-indexed slide 4

    # Remove all old shapes on Slide 4 except background
    for shape in list(slide4.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # 1. Slide Title (36 pt Bold - Matching Slide 5)
    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "🔍 1. 描述（Description：客觀觀察數據與行為記錄）"
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(74, 59, 50) # Matching Slide 5 header color

    # 2. Section 1: 成員與座位配置 (Sub-heading 28 pt Bold + Body 24 pt)
    box1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(248, 250, 252)
    box1.line.color.rgb = RGBColor(226, 232, 240)
    box1.line.width = Pt(1.5)

    tf1 = box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "• 成員與座位配置："
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(29, 78, 216)

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "4位學生圍坐（左前戴紫色口罩黑衣女學生、左後白T男學生、右前米灰T男學生、右後藍衣女學生）。"
    p1_sub.font.size = Pt(24)
    p1_sub.font.color.rgb = RGBColor(30, 41, 59)

    # 3. Section 2: 觀課時間序列紀錄 (Sub-heading 28 pt Bold + 3 Time Sequence Cards in 24~28 pt)
    box2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.85), Inches(11.7), Inches(4.3))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box2.line.color.rgb = RGBColor(203, 213, 225)
    box2.line.width = Pt(1.5)

    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "• 觀課時間序列紀錄："
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(29, 78, 216)
    p2.space_after = Pt(6)

    seq_records = [
        ("【00:00 - 10:00 摸索與確立主筆】", "右前米灰T男生主動拿起彩筆在海報上記錄，左側黑衣女生托腮關注，左後白T男生緊盯海報，形成單一主筆分工。"),
        ("【10:00 - 25:00 觀點對話與資料搜尋】", "進入文本對照，黑衣女生調整筆記本並質疑地形特徵，米灰T男生翻閱課本對照古地圖，引發跨文本對話。"),
        ("【25:00 - 40:00 整合繪製與伸展挑戰】", "整合想法至海報，米灰T男生讓出彩筆給黑衣女生補充繪製，白T男生朗讀關鍵字，展現溫和協同。")
    ]

    for title, desc in seq_records:
        p_title = tf2.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(15, 23, 42)
        p_title.space_before = Pt(4)

        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(23)
        p_desc.font.color.rgb = RGBColor(51, 65, 85)
        p_desc.space_after = Pt(6)

    prs.save(pptx_path)
    print(f"Successfully matched Slide 4 formatting to Slide 5 in {pptx_path}")

if __name__ == "__main__":
    update_slide4_formatting()
