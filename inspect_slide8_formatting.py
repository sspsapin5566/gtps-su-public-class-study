import pptx

prs = pptx.Presentation("output/slides.pptx")

# Inspect Slide 8 formatting (Group 2 Description slide)
slide8 = prs.slides[7] # 0-indexed slide 8
print("=== SLIDE 8 FORMATTING ===")
for shape in slide8.shapes:
    if shape.has_text_frame:
        print(f"Shape: {shape.name}")
        for i, p in enumerate(shape.text_frame.paragraphs):
            size = p.font.size.pt if p.font.size else "Default"
            bold = p.font.bold
            print(f"  P{i} [size={size}, bold={bold}]: {p.text[:90]}...")
