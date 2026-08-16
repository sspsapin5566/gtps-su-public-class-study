import pptx

prs = pptx.Presentation("output/slides.pptx")

# Inspect Slide 5 formatting
slide5 = prs.slides[4] # 0-indexed slide 5
print("=== SLIDE 5 FORMATTING ===")
for shape in slide5.shapes:
    if shape.has_text_frame:
        print(f"Shape: {shape.name}")
        for i, p in enumerate(shape.text_frame.paragraphs):
            size = p.font.size.pt if p.font.size else "Default"
            bold = p.font.bold
            print(f"  P{i} [size={size}, bold={bold}]: {p.text[:80]}...")

print("\n=== SLIDE 4 FORMATTING ===")
slide4 = prs.slides[3] # 0-indexed slide 4
for shape in slide4.shapes:
    if shape.has_text_frame:
        print(f"Shape: {shape.name}")
        for i, p in enumerate(shape.text_frame.paragraphs):
            size = p.font.size.pt if p.font.size else "Default"
            bold = p.font.bold
            print(f"  P{i} [size={size}, bold={bold}]: {p.text[:80]}...")
