import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def update_slide4_to_slide8():
    pptx_path = "output/slides.pptx"
    prs = pptx.Presentation(pptx_path)

    slide4 = prs.slides[3] # 0-indexed slide 4

    # Remove all shapes on Slide 4 except background
    for shape in list(slide4.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # 1. Slide Title (36 pt Bold - Matching Slide 8 Title)
    txBox = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "🔍 1. 描述（Description：客觀觀察數據與行為記錄）"
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(74, 59, 50) # Matching Slide 8 header color (#4A3B32)

    # 2. Main Content Box (Matching Slide 8 font size: 22 pt, with bold title & body distinction)
    content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.6))
    tf_c = content_box.text_frame
    tf_c.word_wrap = True

    # Bullet 1: 成員與座位配置
    p1 = tf_c.paragraphs[0]
    p1.text = "• 成員與座位配置："
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(29, 78, 216) # Dark Blue
    p1.space_after = Pt(2)

    p1_body = tf_c.add_paragraph()
    p1_body.text = "4位學生圍坐（左前戴紫色口罩黑衣女學生、左後白T男學生、右前米灰T男學生、右後藍衣女學生）。"
    p1_body.font.size = Pt(22)
    p1_body.font.bold = False
    p1_body.font.color.rgb = RGBColor(30, 41, 59)
    p1_body.space_after = Pt(12)

    # Bullet 2: 觀課時間序列紀錄
    p2 = tf_c.add_paragraph()
    p2.text = "• 觀課時間序列紀錄："
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(29, 78, 216) # Dark Blue
    p2.space_after = Pt(4)

    seq_records = [
        ("• 【00:00 - 10:00 摸索與確立主筆】：", "右前米灰T男生主動拿起彩筆在大海報上記錄，左側黑衣女生單手托腮觀看，左後白T男生眼睛緊盯海報。形成「單一主筆」的初始分工型態。"),
        ("• 【10:00 - 25:00 觀點對話與資料搜尋】：", "討論進入文本對照階段。左側黑衣女生伸出雙手調整桌面上的筆記本與地圖，向米灰T男生提出質疑：「這裡的地形是不是有變？」米灰T男生隨即翻閱社會課本進行比對。"),
        ("• 【25:00 - 40:00 整合繪製與伸展挑戰】：", "小組將個人想法整合至大海報，米灰T男生讓出彩筆給黑衣女生補充繪製，白T男生負責朗讀課本關鍵字，展現溫和的協同習慣。")
    ]

    for title, desc in seq_records:
        p_seq = tf_c.add_paragraph()
        p_seq.space_after = Pt(8)
        
        # Add Title Run (Bold 22 pt)
        run_t = p_seq.add_run()
        run_t.text = title
        run_t.font.size = Pt(22)
        run_t.font.bold = True
        run_t.font.color.rgb = RGBColor(15, 23, 42)

        # Add Desc Run (Regular 22 pt)
        run_d = p_seq.add_run()
        run_d.text = desc
        run_d.font.size = Pt(22)
        run_d.font.bold = False
        run_d.font.color.rgb = RGBColor(51, 65, 85)

    prs.save(pptx_path)
    print(f"Successfully updated Slide 4 to match Slide 8's 22 pt font size and heading/body distinction in {pptx_path}")

if __name__ == "__main__":
    update_slide4_to_slide8()
