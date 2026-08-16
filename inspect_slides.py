import pptx

prs = pptx.Presentation("output/slides.pptx")
print("Total slides:", len(prs.slides))

for i, slide in enumerate(prs.slides):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            text = shape.text_frame.text.replace("\n", " ")
            if not title:
                title = text[:60]
    print(f"Slide {i+1}: {title}")
